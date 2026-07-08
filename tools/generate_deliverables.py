from __future__ import annotations

import json
import math
import shutil
import zipfile
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt


ROOT = Path(__file__).resolve().parent.parent
DELIVERABLE_DIR = ROOT / "deliverables" / "submission-2026-06-29"
BUILD_DIR = ROOT / "output" / "deliverable-build"
INSTALL_STAGE_DIR = BUILD_DIR / "install-package"
INSTALL_APP_DIR = INSTALL_STAGE_DIR / "Teaching-System"
SOURCE_STAGE_DIR = BUILD_DIR / "source-package"

TODAY = date(2026, 6, 29)
NOW = datetime(2026, 6, 29, 18, 0, 0)
PROJECT_NAME = "基于大模型的软件实训教学结果检查评价与报表系统"
SHORT_NAME = "教学评价系统"
VERSION = "V1.0"
DEFAULT_ACCOUNT = "teacher / 123456"

BACKEND_JAR = ROOT / "target" / "Teaching-System-1.0-SNAPSHOT.jar"
FRONTEND_DIST = ROOT / "frontend" / "dist"
DOCX_RENDER_SCRIPT = ROOT / ".codex-ignore-placeholder"


@dataclass
class PackageStats:
    zip_name: str
    zip_size_mb: float
    file_count: int


def ensure_clean_dir(path: Path) -> None:
    if path.exists():
        shutil.rmtree(path)
    path.mkdir(parents=True, exist_ok=True)


def ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def mb(size: int) -> float:
    return round(size / 1024 / 1024, 2)


def file_count(path: Path) -> int:
    return sum(1 for child in path.rglob("*") if child.is_file())


def copy_tree(src: Path, dst: Path, ignore=None) -> None:
    if dst.exists():
        shutil.rmtree(dst)
    shutil.copytree(src, dst, ignore=ignore)


def remove_paths(base: Path, relative_paths: list[str]) -> None:
    for rel in relative_paths:
        target = base / rel
        if target.is_dir():
            shutil.rmtree(target)
        elif target.exists():
            target.unlink()


def zip_dir(src_dir: Path, zip_path: Path, exclude_prefixes: tuple[str, ...] = ()) -> PackageStats:
    ensure_parent(zip_path)
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for path in sorted(src_dir.rglob("*")):
            if not path.is_file():
                continue
            rel = path.relative_to(src_dir)
            rel_str = rel.as_posix()
            if any(rel_str.startswith(prefix) for prefix in exclude_prefixes):
                continue
            zf.write(path, rel_str)
    return PackageStats(
        zip_name=zip_path.name,
        zip_size_mb=mb(zip_path.stat().st_size),
        file_count=sum(
            1 for p in src_dir.rglob("*")
            if p.is_file() and not any(p.relative_to(src_dir).as_posix().startswith(prefix) for prefix in exclude_prefixes)
        ),
    )


def set_run_font(run, *, name: str = "Microsoft YaHei", size: int | None = None,
                 bold: bool | None = None, color: tuple[int, int, int] | None = None) -> None:
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    run._element.rPr.rFonts.set(qn("w:ascii"), name)
    run._element.rPr.rFonts.set(qn("w:hAnsi"), name)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = RGBColor(*color)


def configure_doc(doc: Document, title: str, subject: str) -> None:
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Calibri"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.15

    for style_name, size, color in [
        ("Heading 1", 16, (46, 116, 181)),
        ("Heading 2", 13, (46, 116, 181)),
        ("Heading 3", 12, (31, 77, 120)),
    ]:
        style = styles[style_name]
        style.font.name = "Calibri"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        style.font.bold = True
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor(*color)

    props = doc.core_properties
    props.author = "Codex"
    props.title = title
    props.subject = subject
    props.comments = "依据仓库现有 README、设计期文档、测试记录与部署材料整理生成。"
    props.version = VERSION
    props.created = NOW
    props.modified = NOW


def add_cover(doc: Document, doc_title: str, subtitle: str, basis: list[str]) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(SHORT_NAME)
    set_run_font(run, size=24, bold=True, color=(11, 37, 69))
    p.space_after = Pt(3)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p2.add_run(doc_title)
    set_run_font(run, size=18, bold=True, color=(46, 116, 181))

    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p3.add_run(subtitle)
    set_run_font(run, size=11, color=(85, 85, 85))

    doc.add_paragraph("")
    table = doc.add_table(rows=4, cols=2)
    table.style = "Table Grid"
    rows = [
        ("项目名称", PROJECT_NAME),
        ("版本", VERSION),
        ("整理日期", TODAY.isoformat()),
        ("内容依据", "；".join(basis)),
    ]
    for i, (k, v) in enumerate(rows):
        table.cell(i, 0).text = k
        table.cell(i, 1).text = v

    doc.add_paragraph("")
    note = doc.add_paragraph()
    note.add_run("说明：").bold = True
    note.add_run("本文件优先复用仓库现有 GitHub 文档内容，对仓库中缺失的成品文档部分进行了补充整理与生成。")
    doc.add_page_break()


def add_bullets(doc: Document, items: list[str], level: int = 0) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        if level:
            p.paragraph_format.left_indent = Inches(0.25 * level)
        p.add_run(item)


def add_numbers(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Number")
        p.add_run(item)


def add_table(doc: Document, headers: list[str], rows: list[list[str]]) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    for row in rows:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            cells[idx].text = value


def add_sources_section(doc: Document, sources: list[str]) -> None:
    doc.add_heading("资料来源", level=1)
    add_bullets(doc, sources)


def create_requirements_doc(path: Path) -> None:
    doc = Document()
    configure_doc(doc, "软件功能需求分析文档", "需求分析")
    add_cover(
        doc,
        "软件功能需求分析文档",
        "基于仓库现有需求与设计期材料整理",
        [
            "README.md",
            "docs/设计期文档/01-项目概述与需求分析.md",
            "docs/设计期文档/02-业务流程设计.md",
            "docs/设计期文档/12-角色与权限说明.md",
            "docs/设计期文档/13-阶段计划与范围边界.md",
        ],
    )

    doc.add_heading("1. 文档目的", level=1)
    doc.add_paragraph(
        "本需求分析文档用于说明教学评价系统当前阶段要解决的问题、服务对象、功能边界、核心业务流程和验收基线，"
        "作为后续设计、开发、测试和演示交付的统一需求依据。"
    )

    doc.add_heading("2. 项目背景与目标", level=1)
    doc.add_paragraph(
        "在软件实训教学过程中，教师需要检查学生提交的代码、实验报告、截图和压缩包材料，"
        "传统纯人工评价方式在批量场景下存在效率低、标准不统一、统计导出麻烦等问题。"
        "本系统通过“作品管理 + AI 辅助评价 + 教师复核 + 统计导出”的闭环，帮助教师提高教学评价效率。"
    )
    add_bullets(doc, [
        "降低教师批量检查学生作品的重复劳动。",
        "提供可复核的 AI 建议分数、问题清单和评语草稿。",
        "形成统一的教师最终成绩、统计摘要和 Excel 导出结果。",
        "保留真实大模型、真实文件上传和多角色扩展的接口空间。",
    ])

    doc.add_heading("3. 用户对象与范围", level=1)
    add_table(doc, ["项", "当前阶段定义"], [
        ["主要使用者", "教师"],
        ["登录角色", "固定教师账号，当前不开放学生端与管理员端"],
        ["评价对象", "学生提交的代码、文档、截图、压缩包等作品材料"],
        ["当前范围", "教师登录、班级管理、作业管理、作品提交、AI 评价、教师复核、统计、Excel 导出"],
        ["非本阶段范围", "完整多角色权限、学生自主提交门户、复杂审批链、生产级多租户部署"],
    ])

    doc.add_heading("4. 核心业务流程", level=1)
    add_numbers(doc, [
        "教师登录系统，进入主工作台。",
        "维护班级、学生与作业基础信息。",
        "登记或上传学生作品，形成提交记录。",
        "触发 AI 自动评价，生成建议分数、问题分析和评语。",
        "教师查看 AI 结果并进行最终复核与确认。",
        "系统汇总学生数、提交数、AI 已评价数、教师已确认数和平均分。",
        "导出 Excel 成绩表，完成教学评价闭环。",
    ])

    doc.add_heading("5. 功能需求", level=1)
    add_table(doc, ["模块", "功能要求", "当前实现状态"], [
        ["登录认证", "教师账号登录、会话保持、退出登录", "已实现"],
        ["班级管理", "班级查看、新增、编辑、关联学生", "已实现"],
        ["作业管理", "作业新增、修改、删除、按班级筛选", "已实现"],
        ["作品提交", "提交记录新增、真实文件上传、文件明细记录", "已实现"],
        ["文件预览/下载", "图片/视频在线预览、主文件与明细文件下载", "已实现"],
        ["AI 评价", "调用真实 AI 或降级到演示评分，返回分数与评语", "已实现"],
        ["教师复核", "教师最终评分、评语保存、状态迁移", "已实现"],
        ["统计分析", "仪表盘统计摘要、评分分布和趋势展示", "已实现"],
        ["成绩导出", "Excel 文件流导出", "已实现"],
        ["审计与安全", "登录校验、防跨站请求头、速率限制、审计日志", "已实现"],
    ])

    doc.add_heading("6. 非功能需求", level=1)
    add_bullets(doc, [
        "易用性：页面入口清晰，教师能在一个工作台完成主要业务操作。",
        "可维护性：前后端分离，后端 API 与 Python AI 服务职责明确。",
        "可扩展性：保留真实大模型、多文件上传、更多文件类型预处理的扩展接口。",
        "安全性：关键 API 需要登录，会话使用 HttpOnly Cookie，写操作要求防跨站请求头。",
        "可部署性：支持本地一键启动、Docker Compose 部署和 LoongArch 麒麟原生演示。",
    ])

    doc.add_heading("7. 数据与接口需求概要", level=1)
    doc.add_paragraph("系统核心数据对象包括班级、学生、作业、作品提交、提交文件和评价结果。")
    add_bullets(doc, [
        "班级与学生形成基础教学组织结构。",
        "作业与班级支持一对多及多班级受理关系。",
        "作品提交与提交文件支持真实文件元数据和主文件快照。",
        "评价结果记录 AI 分数、问题、评语、教师分数、教师评语和状态。",
        "详细接口定义以 Swagger / OpenAPI 为准，需求文档只保留总体要求和业务含义。",
    ])

    doc.add_heading("8. 当前验收基线", level=1)
    add_bullets(doc, [
        "主线代码基于 Vue 3 前端、Spring Boot 后端和 Python AI 服务。",
        "演示基线可完成“登录 -> 作业审批 -> 附件预览/下载 -> AI 评价 -> 教师复核 -> 统计/Excel 导出”。",
        "在真实 AI 服务不可用时，系统允许自动降级到演示评分，以保证教学演示闭环。",
        "LoongArch 麒麟环境优先使用原生启动脚本完成最终验收演示。",
    ])

    add_sources_section(doc, [
        "README.md",
        "docs/设计期文档/01-项目概述与需求分析.md",
        "docs/设计期文档/02-业务流程设计.md",
        "docs/设计期文档/12-角色与权限说明.md",
        "docs/设计期文档/13-阶段计划与范围边界.md",
    ])
    doc.save(path)


def create_design_doc(path: Path) -> None:
    doc = Document()
    configure_doc(doc, "软件功能设计文档", "系统设计")
    add_cover(
        doc,
        "软件功能设计文档",
        "基于仓库现有设计期文档与主线代码整理",
        [
            "README.md",
            "docs/设计期文档/00-15 系列文档",
            "src/main/java/com/teachingeval/",
            "frontend/src/",
            "ai-service/",
        ],
    )

    doc.add_heading("1. 设计目标", level=1)
    doc.add_paragraph(
        "本设计文档用于说明教学评价系统的总体架构、模块划分、数据设计、接口设计、页面设计和部署设计，"
        "保证需求分析、主线代码和验收演示材料在设计口径上保持一致。"
    )

    doc.add_heading("2. 总体技术架构", level=1)
    add_table(doc, ["层次", "技术栈", "职责"], [
        ["前端展示层", "Vue 3 + Vite + TypeScript + SCSS", "负责教师端 SPA 页面、状态管理、交互体验和文件预览入口"],
        ["业务服务层", "Spring Boot 3.3 + Spring Data JPA + H2/MySQL", "负责认证、班级、作业、提交、评价、统计和导出接口"],
        ["AI 处理层", "Python FastAPI + DocxConv + Evaluator + OCR/Archive 模块", "负责文件预处理、真实 AI 评分与健康检查"],
        ["数据存储层", "H2（开发）、MySQL（正式）", "负责教学评价主数据存储"],
    ])

    doc.add_heading("3. 模块设计", level=1)
    add_table(doc, ["模块", "主要类/目录", "设计说明"], [
        ["认证模块", "AuthController / AuthService / ApiSecurityFilter", "使用教师账号和 HttpOnly Cookie 管理会话"],
        ["班级与学生模块", "TeachingClassController / StudentController / Service / Repository", "维护班级与学生基础数据"],
        ["作业管理模块", "AssignmentController / AssignmentService", "管理作业及作业与班级的关联关系"],
        ["作品提交模块", "SubmissionController / SubmissionService", "处理提交记录、真实文件上传、文件下载和预览"],
        ["评价模块", "EvaluationController / EvaluationService / DelegatingAIService", "根据配置调用真实 AI 或演示评分"],
        ["统计导出模块", "StatisticsController / ExportController / ExportService", "输出统计摘要与 Excel 文件流"],
        ["AI 服务模块", "ai-service/DocxConv / Evaluator / ScreenshotProc / ArchiveProc", "提供文件预处理与 AI 能力"],
    ])

    doc.add_heading("4. 页面与交互设计", level=1)
    add_table(doc, ["页面/路由", "作用", "主要交互"], [
        ["/login", "教师登录页", "账号密码登录、登录失败提示"],
        ["/dashboard", "仪表盘", "统计摘要、趋势信息、状态概览"],
        ["/classes", "班级管理", "班级列表、学生花名册、编辑与删除"],
        ["/assignments", "作业管理", "作业列表、筛选、编辑、Excel 导出"],
        ["/review", "作业审批页", "附件区、AI 评价、教师复核、草稿恢复"],
        ["/file-preview", "文件预览页", "文档、图片、视频预览"],
        ["/forbidden", "未登录/无权限页", "登录过期或无权限提示"],
    ])

    doc.add_heading("5. 数据库设计摘要", level=1)
    add_table(doc, ["数据表", "作用", "关键字段"], [
        ["teaching_class", "班级信息", "id, name, grade, description"],
        ["student", "学生信息", "student_no, name, class_id, class_name"],
        ["assignment", "作业信息", "title, work_type, class_id, due_at"],
        ["assignment_class", "作业受理班级关联", "assignment_id, class_id"],
        ["submission", "作品提交记录", "student_id, assignment_id, file_name, preprocess_status"],
        ["submission_file", "提交文件明细", "submission_id, file_name, file_path, primary_file"],
        ["evaluation", "评价结果", "submission_id, ai_score, teacher_score, status"],
    ])

    doc.add_heading("6. 接口设计摘要", level=1)
    add_bullets(doc, [
        "所有 REST API 统一使用 /api 前缀。",
        "资源命名遵循名词复数与短横线规范，例如 /api/submissions、/api/submissions/{id}/teacher-review。",
        "接口详细路径、请求参数、返回字段以 Swagger UI 为权威说明。",
        "文件流、Excel 导出和 SSE 流式评分按非普通 JSON 响应处理。",
    ])
    add_table(doc, ["接口模块", "代表接口", "说明"], [
        ["登录认证", "POST /api/auth/login, POST /api/auth/logout, GET /api/auth/me", "教师登录、退出、会话查询"],
        ["班级/学生", "GET /api/classes, GET/POST /api/students", "班级和学生基础数据维护"],
        ["作业", "GET/POST/PUT/DELETE /api/assignments", "作业全生命周期管理"],
        ["提交", "GET/POST /api/submissions, POST /api/submissions/upload", "提交记录与真实文件上传"],
        ["评价", "POST /api/submissions/{id}/evaluate, POST /api/submissions/{id}/teacher-review", "AI 评价与教师复核"],
        ["统计导出", "GET /api/statistics/summary, POST /api/export/excel", "统计摘要与 Excel 导出"],
    ])

    doc.add_heading("7. AI 服务设计", level=1)
    add_bullets(doc, [
        "Java 侧通过 DelegatingAIService 对真实 AI 和 FakeAIService 做统一封装。",
        "当 app.ai.real.enabled=true 时，后端优先调用 Python /api/evaluate/real。",
        "当真实 AI 不可用或无真实文件时，系统允许降级为演示评分，保证业务闭环不被外部依赖阻断。",
        "Python /api/preprocess 负责文档解析、OCR、压缩包解压和文本提取。",
    ])

    doc.add_heading("8. 安全与审计设计", level=1)
    add_bullets(doc, [
        "ApiSecurityFilter 校验登录状态和防跨站请求头。",
        "AuthController 通过 HttpOnly Cookie 管理会话令牌。",
        "RateLimitService 对 AI 评价接口做分钟级调用限制。",
        "AuditLogFilter / AuditLogController 记录关键操作轨迹。",
    ])

    doc.add_heading("9. 部署设计", level=1)
    add_table(doc, ["部署方式", "适用场景", "说明"], [
        ["本地一键启动", "开发调试", "使用 start.sh / start.bat 启动三端服务"],
        ["LoongArch 麒麟原生演示", "课程验收演示", "使用 start-kylin-loongarch-demo.sh，优先保证演示稳定性"],
        ["Docker Compose", "普通 Linux / 容器部署", "保留 ai-service、backend、frontend 三服务编排能力"],
    ])

    add_sources_section(doc, [
        "README.md",
        "docs/设计期文档/00-设计期文档目录.md",
        "docs/设计期文档/02-业务流程设计.md",
        "docs/设计期文档/03-项目结构.md",
        "docs/设计期文档/04-数据库设计.md",
        "docs/设计期文档/05-接口设计总体方针.md",
        "docs/设计期文档/06-页面设计.md",
        "docs/设计期文档/07-组件交互设计.md",
        "docs/设计期文档/10-数据流说明.md",
        "docs/设计期文档/14-Excel导出设计.md",
        "docs/设计期文档/15-AI服务设计.md",
    ])
    doc.save(path)


def create_manual_doc(path: Path) -> None:
    doc = Document()
    configure_doc(doc, "软件产品说明书", "产品说明书")
    add_cover(
        doc,
        "软件产品说明书",
        "教师端操作说明与常见问题",
        ["README.md", "frontend/src/views/", "docs/演示文档/麒麟LoongArch虚拟机部署指南.md"],
    )

    doc.add_heading("1. 产品概述", level=1)
    doc.add_paragraph(
        f"{SHORT_NAME}面向教师使用，围绕班级、作业、作品提交、AI 评价、教师复核、统计分析和 Excel 导出，"
        "提供一套完整的教学评价闭环。当前主线为前后端分离结构，支持本地演示、Docker 部署和 LoongArch 麒麟环境演示。"
    )

    doc.add_heading("2. 使用环境", level=1)
    add_table(doc, ["项目", "要求"], [
        ["浏览器", "Chrome、Edge 或其他现代 Chromium 浏览器"],
        ["后端运行环境", "Java 17、Maven 3.9+"],
        ["前端运行环境", "Node.js 24、npm 11+"],
        ["AI 服务（可选）", "Python 3.13 虚拟环境或 Docker 容器"],
        ["推荐演示账号", DEFAULT_ACCOUNT],
    ])

    doc.add_heading("3. 登录与退出", level=1)
    add_numbers(doc, [
        "打开前端页面地址，例如 http://localhost:5173。",
        "输入教师账号与密码。",
        "登录成功后进入仪表盘，系统会在后端写入 HttpOnly 会话 Cookie。",
        "需要退出时，可通过页面退出按钮或调用 /api/auth/logout 清除会话。",
    ])

    doc.add_heading("4. 主要功能说明", level=1)
    doc.add_heading("4.1 仪表盘", level=2)
    add_bullets(doc, [
        "查看学生总数、作品提交数、AI 已评价数、教师已确认数和平均分。",
        "快速感知当前教学评价进度。",
    ])

    doc.add_heading("4.2 班级管理", level=2)
    add_bullets(doc, [
        "查看班级列表和班级基础信息。",
        "新增、编辑或删除班级与学生花名册信息。",
        "为作业、提交和统计提供基础数据。"
    ])

    doc.add_heading("4.3 作业管理", level=2)
    add_bullets(doc, [
        "新增作业，设置作业标题、类型、班级和截止时间。",
        "对现有作业进行查询、筛选、编辑与删除。",
        "在作业管理页面直接触发 Excel 导出。"
    ])

    doc.add_heading("4.4 作品提交与文件处理", level=2)
    add_bullets(doc, [
        "录入作品标题、作品类型、备注等元数据。",
        "支持真实文件上传，记录主文件和文件明细。",
        "支持图片、视频等文件的在线预览与原文件下载。"
    ])

    doc.add_heading("4.5 AI 评价与教师复核", level=2)
    add_bullets(doc, [
        "教师在审批页触发 AI 评价，系统返回建议分数、问题清单和评语。",
        "当真实 AI 不可用时，系统会自动降级到演示评分模式。",
        "教师可以修改最终分数和评语，并保存为教师确认结果。"
    ])

    doc.add_heading("4.6 统计与导出", level=2)
    add_bullets(doc, [
        "系统根据作品与评价状态实时计算统计摘要。",
        "可导出 .xlsx 成绩文件，供教学归档和后续统计使用。"
    ])

    doc.add_heading("5. 推荐操作流程", level=1)
    add_numbers(doc, [
        "登录系统。",
        "确认班级与学生基础数据。",
        "创建作业或检查作业信息。",
        "上传/登记学生作品。",
        "在作业审批页执行 AI 评价。",
        "教师完成最终复核并保存。",
        "查看统计结果并导出 Excel。"
    ])

    doc.add_heading("6. 常见问题", level=1)
    add_table(doc, ["问题", "处理建议"], [
        ["登录后接口返回未授权", "确认浏览器已接收 auth_token Cookie，且请求未被跨站限制。"],
        ["AI 评价失败", "确认 Python AI 服务是否已启动；若未启动，演示环境可接受自动降级为模拟评分。"],
        ["文件预览失败", "确认文件类型是否属于当前支持的图片/视频预览范围，或改为直接下载查看。"],
        ["Excel 导出无响应", "确认后端服务正常，浏览器允许文件下载，且评价数据至少存在已确认结果。"],
    ])

    doc.add_heading("7. 使用注意事项", level=1)
    add_bullets(doc, [
        "当前系统面向教师端演示与课程验收，不等同于完整生产环境教学平台。",
        "真实 AI、OCR 和部分格式预处理依赖外部环境，可在演示时采用降级方案。",
        "在 LoongArch 麒麟验收环境中，优先采用仓库内原生启动方案。"
    ])

    add_sources_section(doc, [
        "README.md",
        "frontend/src/views/LoginPage.vue",
        "frontend/src/views/DashboardPage.vue",
        "frontend/src/views/ClassesPage.vue",
        "frontend/src/views/AssignmentsPage.vue",
        "frontend/src/views/ReviewPage.vue",
        "docs/演示文档/麒麟LoongArch虚拟机部署指南.md",
    ])
    doc.save(path)


def create_test_doc(path: Path) -> None:
    doc = Document()
    configure_doc(doc, "软件功能测试报告", "测试报告")
    add_cover(
        doc,
        "软件功能测试报告",
        "基于 2026-06-29 本地实际测试结果整理",
        [
            "mvn test",
            "npm test",
            "npm run build",
            "mvn -DskipTests package",
            "docs/开发期文档/测试记录/测试记录.md",
        ],
    )

    doc.add_heading("1. 测试目的", level=1)
    doc.add_paragraph(
        "验证教学评价系统当前主线代码在 2026-06-29 的本地环境下具备可编译、可测试、可构建和可演示的基本交付能力，"
        "并对关键业务闭环进行结果归档。"
    )

    doc.add_heading("2. 测试环境", level=1)
    add_table(doc, ["项", "结果"], [
        ["测试日期", "2026-06-29"],
        ["操作系统", "macOS / Darwin 24.6.0 / arm64"],
        ["Java", "OpenJDK 17.0.18 LTS"],
        ["Maven", "3.9.16"],
        ["Node.js", "24.14.0"],
        ["npm", "11.9.0"],
        ["数据库", "H2 内存数据库（测试场景）"],
    ])

    doc.add_heading("3. 自动化测试结果", level=1)
    add_table(doc, ["测试类别", "执行命令", "结果", "说明"], [
        ["后端单元/集成测试", "mvn test", "通过（24/24）", "TeachingSystemSecurityTest 9 项，TeachingSystemFlowTest 15 项，BUILD SUCCESS"],
        ["前端单元测试", "npm test", "通过（2/2）", "Vitest 运行通过，1 个测试文件、2 个测试用例"],
        ["前端生产构建", "npm run build", "通过", "Vite 构建成功，存在大包体提示但不影响构建成功"],
        ["后端打包验证", "mvn -DskipTests package", "通过", "成功生成可执行 Spring Boot JAR"],
    ])

    doc.add_heading("4. 关键测试说明", level=1)
    add_bullets(doc, [
        "后端测试覆盖登录保护、会话 Cookie、写操作请求头校验、文件下载权限、AI 评价频率限制等安全场景。",
        "后端测试覆盖学生分页、班级/作业关系、提交流程、AI 评价、教师复核、统计汇总和真实文件上传流程。",
        "前端测试当前已覆盖组件级基础测试，当前重点在于保证主线页面构建和联调链路稳定。",
    ])

    doc.add_heading("5. 测试中观察到的现象", level=1)
    add_bullets(doc, [
        "在执行后端测试时，日志中出现“真实 AI 服务调用失败，已降级为模拟评价”的警告，这是系统设计允许的降级路径，最终测试结果仍为通过。",
        "前端生产构建存在若干超过 500kB 的 chunk 告警，属于体积优化建议，不影响构建产物生成。",
    ])

    doc.add_heading("6. 构建产物核验", level=1)
    add_table(doc, ["产物", "位置", "结果"], [
        ["后端 JAR", "target/Teaching-System-1.0-SNAPSHOT.jar", "已生成"],
        ["兼容 JAR 副本", "target/app.jar", "已存在"],
        ["前端静态构建目录", "frontend/dist/", "已生成"],
        ["测试报告 XML", "target/surefire-reports/*.xml", "已生成"],
    ])

    doc.add_heading("7. 结论", level=1)
    doc.add_paragraph(
        "截至 2026-06-29，本项目主线代码已完成自动化测试、前端构建和后端打包验证，"
        "具备形成课程验收交付物的基础条件。建议在最终 LoongArch 麒麟演示环境中再执行一次原生启动与录屏复测。"
    )

    add_sources_section(doc, [
        "target/surefire-reports/TEST-com.teachingeval.TeachingSystemFlowTest.xml",
        "target/surefire-reports/TEST-com.teachingeval.TeachingSystemSecurityTest.xml",
        "docs/开发期文档/测试记录/测试记录.md",
    ])
    doc.save(path)


def create_deployment_doc(path: Path, install_stats: PackageStats, source_stats: PackageStats) -> None:
    doc = Document()
    configure_doc(doc, "软件安装包及部署文档", "部署文档")
    add_cover(
        doc,
        "软件安装包及部署文档",
        "包含安装包内容说明、部署方式和验收建议",
        [
            "README.md",
            "start.sh / start.bat / start-kylin-loongarch-demo.sh",
            "docker-compose.yml",
            "docs/演示文档/麒麟LoongArch虚拟机部署指南.md",
        ],
    )

    doc.add_heading("1. 交付包概述", level=1)
    add_table(doc, ["交付内容", "文件名", "说明"], [
        ["安装包", install_stats.zip_name, f"共 {install_stats.file_count} 个文件，压缩后约 {install_stats.zip_size_mb} MB"],
        ["源码包", source_stats.zip_name, f"共 {source_stats.file_count} 个文件，压缩后约 {source_stats.zip_size_mb} MB"],
        ["部署说明", path.name, "说明不同环境下的部署方式与检查步骤"],
    ])

    doc.add_heading("2. 安装包内容", level=1)
    add_bullets(doc, [
        "后端可执行 JAR：Teaching-System-1.0-SNAPSHOT.jar",
        "前端构建产物：frontend/dist/",
        "AI 服务源码与依赖说明：ai-service/",
        "Docker 与原生部署脚本：docker-compose.yml、Dockerfile、start.sh、start.bat、start-kylin-loongarch-demo.sh",
        "数据库脚本：sql/schema.sql",
    ])

    doc.add_heading("3. 部署方式一：本地开发/演示启动", level=1)
    add_numbers(doc, [
        "准备 Java 17、Maven、Node.js 和 npm 环境。",
        "执行 Linux/macOS 启动脚本 ./start.sh，或 Windows 启动脚本 start.bat。",
        "如需 AI 服务完整能力，按 README 中步骤安装 Python 依赖并启动 ai-service。",
        "访问前端地址 http://localhost:5173，使用教师账号登录。",
    ])

    doc.add_heading("4. 部署方式二：LoongArch 麒麟原生演示", level=1)
    add_bullets(doc, [
        "推荐用于课程最终验收，优先保证“环境合规 + 功能闭环”而不是依赖复杂容器镜像。",
        "在 LoongArch 麒麟系统中执行 chmod +x start-kylin-loongarch-demo.sh && ./start-kylin-loongarch-demo.sh。",
        "演示前建议展示 uname -m 与 /etc/os-release，证明环境为 loongarch64 + 麒麟系统。",
    ])

    doc.add_heading("5. 部署方式三：Docker Compose", level=1)
    add_numbers(doc, [
        "执行 mvn -DskipTests package 生成后端 JAR。",
        "执行 docker compose up --build -d 启动 ai-service、backend、frontend 三个容器。",
        "通过 /api/health、Swagger UI 和前端首页检查服务状态。",
    ])

    doc.add_heading("6. 健康检查与默认访问地址", level=1)
    add_table(doc, ["服务", "默认地址", "说明"], [
        ["前端页面", "http://localhost:5173", "教师端 SPA"],
        ["后端健康检查", "http://localhost:8080/api/health", "检查后端是否启动"],
        ["Swagger UI", "http://localhost:8080/swagger-ui/index.html", "查看 API 文档"],
        ["AI 服务健康检查（可选）", "http://localhost:8000/api/health", "检查 Python AI 服务"],
    ])

    doc.add_heading("7. 默认账号与运行说明", level=1)
    add_bullets(doc, [
        f"默认教师账号：{DEFAULT_ACCOUNT}",
        "当前演示基线允许在真实 AI 不可用时自动降级为演示评分。",
        "LoongArch 演示环境中若 Docker 镜像或 OCR 依赖不稳定，建议优先采用原生启动脚本。",
    ])

    doc.add_heading("8. 常见问题处理", level=1)
    add_table(doc, ["问题", "处理方式"], [
        ["端口 8080 或 5173 被占用", "先释放端口，再重新执行启动脚本。"],
        ["后端启动失败", "检查 JDK 版本、Maven 依赖下载情况和 logs 中的启动日志。"],
        ["前端启动失败", "检查 npm install 是否完成，必要时切换镜像源。"],
        ["LoongArch 环境 Docker 构建失败", "切换到原生启动方案进行演示。"],
    ])

    add_sources_section(doc, [
        "README.md",
        "start.sh",
        "start.bat",
        "start-kylin-loongarch-demo.sh",
        "docker-compose.yml",
        "docs/演示文档/麒麟LoongArch虚拟机部署指南.md",
    ])
    doc.save(path)


def add_textbox(slide, left, top, width, height, text, *, font_size=20, bold=False,
                color=(20, 33, 61), align=PP_ALIGN.LEFT, font_name="Microsoft YaHei",
                fill=None, line=None, margin=0.08):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGBColor(*fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = PptxRGBColor(*line)
    shape.adjustments[0] = 0.08
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = PptxInches(margin)
    tf.margin_right = PptxInches(margin)
    tf.margin_top = PptxInches(margin)
    tf.margin_bottom = PptxInches(margin)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = PptxPt(font_size)
    font.bold = bold
    font.color.rgb = PptxRGBColor(*color)
    return shape


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    tx = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(0.5), PptxInches(8.8), PptxInches(0.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = PptxPt(24)
    run.font.bold = True
    run.font.color.rgb = PptxRGBColor(15, 45, 74)
    if subtitle:
        st = slide.shapes.add_textbox(PptxInches(0.82), PptxInches(1.1), PptxInches(8.6), PptxInches(0.5))
        sp = st.text_frame.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.name = "Microsoft YaHei"
        srun.font.size = PptxPt(10.5)
        srun.font.color.rgb = PptxRGBColor(96, 96, 96)


def decorate_slide(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PptxRGBColor(247, 244, 236)
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(13.33), PptxInches(0.18)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = PptxRGBColor(15, 45, 74)
    band.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptxInches(11.65), PptxInches(0.18), PptxInches(1.68), PptxInches(7.32)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PptxRGBColor(230, 164, 74)
    accent.fill.transparency = 0.82
    accent.line.fill.background()


def add_bullet_box(slide, left, top, width, height, title, items):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = PptxRGBColor(255, 255, 255)
    box.line.color.rgb = PptxRGBColor(222, 222, 222)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = PptxInches(0.18)
    tf.margin_right = PptxInches(0.12)
    tf.margin_top = PptxInches(0.12)
    tf.margin_bottom = PptxInches(0.12)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.size = PptxPt(16)
    r.font.bold = True
    r.font.color.rgb = PptxRGBColor(15, 45, 74)

    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = PptxPt(10.5)
        p.font.color.rgb = PptxRGBColor(60, 60, 60)
    return box


def create_ppt(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    decorate_slide(slide)
    add_title(slide, "教学评价系统", "软件功能演示文档")
    add_textbox(
        slide, PptxInches(0.85), PptxInches(1.9), PptxInches(6.9), PptxInches(1.5),
        PROJECT_NAME, font_size=20, bold=True, fill=(255, 255, 255), line=(232, 232, 232)
    )
    add_textbox(
        slide, PptxInches(0.85), PptxInches(3.75), PptxInches(4.5), PptxInches(1.0),
        "目标：用 AI 辅助教师完成作品检查、复核评分、统计和导出。", font_size=12.5,
        fill=(255, 250, 242), line=(235, 196, 125)
    )
    add_textbox(
        slide, PptxInches(0.85), PptxInches(5.1), PptxInches(4.5), PptxInches(0.72),
        f"默认账号：{DEFAULT_ACCOUNT}", font_size=12, bold=True, fill=(15, 45, 74),
        color=(255, 255, 255)
    )

    slide = prs.slides.add_slide(blank)
    decorate_slide(slide)
    add_title(slide, "项目背景与目标", "为什么需要这套系统")
    add_bullet_box(slide, PptxInches(0.8), PptxInches(1.6), PptxInches(3.7), PptxInches(4.7), "痛点", [
        "学生作品材料多，教师纯人工检查耗时大",
        "评价标准不易统一，成绩统计与导出繁琐",
        "文档、代码、截图、压缩包需要分散检查",
    ])
    add_bullet_box(slide, PptxInches(4.8), PptxInches(1.6), PptxInches(3.7), PptxInches(4.7), "目标", [
        "把作品管理、AI 评价、复核和导出串成闭环",
        "让教师更快完成教学评价与结果归档",
        "为真实大模型和更多文件类型扩展预留接口",
    ])
    add_bullet_box(slide, PptxInches(8.8), PptxInches(1.6), PptxInches(3.1), PptxInches(4.7), "当前范围", [
        "教师端登录",
        "班级/作业/提交管理",
        "AI 评价与教师复核",
        "统计摘要与 Excel 导出",
    ])

    slide = prs.slides.add_slide(blank)
    decorate_slide(slide)
    add_title(slide, "系统架构", "前端、后端、AI 服务与数据层")
    boxes = [
        (1.0, 2.0, 2.3, 1.0, "Vue 3 前端", (255, 255, 255), (15, 45, 74)),
        (3.9, 2.0, 2.6, 1.0, "Spring Boot 后端", (255, 255, 255), (15, 45, 74)),
        (7.1, 2.0, 2.4, 1.0, "Python AI 服务", (255, 255, 255), (15, 45, 74)),
        (10.1, 2.0, 1.9, 1.0, "H2 / MySQL", (255, 255, 255), (15, 45, 74)),
    ]
    shapes = []
    for left, top, width, height, text, fill, line in boxes:
        shapes.append(add_textbox(
            slide, PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height),
            text, font_size=16, bold=True, fill=fill, line=line, align=PP_ALIGN.CENTER
        ))
    for idx in range(len(shapes) - 1):
        start = shapes[idx]
        end = shapes[idx + 1]
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start.left + start.width,
            start.top + start.height // 2,
            end.left,
            end.top + end.height // 2,
        )
        line.line.color.rgb = PptxRGBColor(230, 164, 74)
        line.line.width = PptxPt(2)
    add_bullet_box(slide, PptxInches(1.0), PptxInches(3.55), PptxInches(4.1), PptxInches(2.0), "前端职责", [
        "教师端路由与页面交互",
        "文件预览、审批、导出入口",
        "MagicBar / 三面板工作流",
    ])
    add_bullet_box(slide, PptxInches(5.3), PptxInches(3.55), PptxInches(4.1), PptxInches(2.0), "后端职责", [
        "认证、权限、业务服务",
        "统计汇总与 Excel 导出",
        "真实 AI / 演示评分统一封装",
    ])
    add_bullet_box(slide, PptxInches(9.55), PptxInches(3.55), PptxInches(2.45), PptxInches(2.0), "AI 服务职责", [
        "文档解析",
        "OCR/压缩包处理",
        "真实 AI 评分",
    ])

    slide = prs.slides.add_slide(blank)
    decorate_slide(slide)
    add_title(slide, "核心业务流程", "从作品接收到成绩导出")
    flow_texts = ["教师登录", "班级/作业管理", "作品提交", "AI 自动评价", "教师复核", "统计导出"]
    x = 0.9
    for i, text in enumerate(flow_texts, start=1):
        shape = add_textbox(
            slide, PptxInches(x), PptxInches(2.55), PptxInches(1.8), PptxInches(1.0),
            f"{i}. {text}", font_size=14, bold=True, fill=(255, 255, 255), line=(15, 45, 74),
            align=PP_ALIGN.CENTER
        )
        if i < len(flow_texts):
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                shape.left + shape.width,
                shape.top + shape.height // 2,
                shape.left + shape.width + PptxInches(0.28),
                shape.top + shape.height // 2,
            )
            line.line.color.rgb = PptxRGBColor(230, 164, 74)
            line.line.width = PptxPt(2.4)
        x += 2.0
    add_textbox(
        slide, PptxInches(1.0), PptxInches(4.35), PptxInches(10.8), PptxInches(1.2),
        "闭环特点：真实 AI 可用时走 Python 评分链路，不可用时自动降级到演示评分，保证课程验收主流程稳定可演示。",
        font_size=13, fill=(255, 250, 242), line=(235, 196, 125)
    )

    slide = prs.slides.add_slide(blank)
    decorate_slide(slide)
    add_title(slide, "功能模块一览", "当前主线页面与模块分工")
    add_bullet_box(slide, PptxInches(0.8), PptxInches(1.6), PptxInches(2.7), PptxInches(4.6), "仪表盘", [
        "展示学生、提交、AI、复核统计",
        "作为教师工作入口页",
    ])
    add_bullet_box(slide, PptxInches(3.7), PptxInches(1.6), PptxInches(2.7), PptxInches(4.6), "班级管理", [
        "班级列表与花名册",
        "学生数据维护",
    ])
    add_bullet_box(slide, PptxInches(6.6), PptxInches(1.6), PptxInches(2.7), PptxInches(4.6), "作业管理", [
        "作业新增、编辑、筛选",
        "Excel 导出入口",
    ])
    add_bullet_box(slide, PptxInches(9.5), PptxInches(1.6), PptxInches(2.0), PptxInches(4.6), "作业审批", [
        "附件预览",
        "AI 评价",
        "教师复核",
        "状态迁移",
    ])

    slide = prs.slides.add_slide(blank)
    decorate_slide(slide)
    add_title(slide, "测试与质量状态", "2026-06-29 本地验证结果")
    metrics = [
        ("后端测试", "24 / 24", "Security 9 + Flow 15"),
        ("前端测试", "2 / 2", "Vitest 通过"),
        ("前端构建", "Success", "Vite 生产构建完成"),
        ("后端打包", "Success", "JAR 产物已生成"),
    ]
    lefts = [0.9, 3.55, 6.2, 8.85]
    for left, (title, value, note) in zip(lefts, metrics):
        add_textbox(
            slide, PptxInches(left), PptxInches(2.0), PptxInches(2.2), PptxInches(2.2),
            f"{title}\n{value}\n{note}", font_size=15, bold=True, align=PP_ALIGN.CENTER,
            fill=(255, 255, 255), line=(15, 45, 74)
        )
    add_textbox(
        slide, PptxInches(0.9), PptxInches(4.75), PptxInches(10.8), PptxInches(1.0),
        "说明：后端测试日志中出现真实 AI 不可用时的降级告警，但属于设计允许的兜底路径，最终测试结果均为通过。",
        font_size=12.5, fill=(255, 250, 242), line=(235, 196, 125)
    )

    slide = prs.slides.add_slide(blank)
    decorate_slide(slide)
    add_title(slide, "部署方式", "开发、验收与容器化三种路径")
    add_bullet_box(slide, PptxInches(0.9), PptxInches(1.7), PptxInches(3.3), PptxInches(4.6), "本地一键启动", [
        "start.sh / start.bat",
        "适合开发调试",
        "前后端与 AI 服务可分开启动",
    ])
    add_bullet_box(slide, PptxInches(4.55), PptxInches(1.7), PptxInches(3.3), PptxInches(4.6), "LoongArch 麒麟原生演示", [
        "start-kylin-loongarch-demo.sh",
        "优先用于课程最终验收",
        "更适合证明环境合规",
    ])
    add_bullet_box(slide, PptxInches(8.2), PptxInches(1.7), PptxInches(3.3), PptxInches(4.6), "Docker Compose", [
        "ai-service + backend + frontend",
        "适合普通 Linux 容器部署",
        "LoongArch 环境可作为备选",
    ])

    slide = prs.slides.add_slide(blank)
    decorate_slide(slide)
    add_title(slide, "本次交付物", "按验收清单整理完成")
    add_bullet_box(slide, PptxInches(0.9), PptxInches(1.6), PptxInches(10.5), PptxInches(4.9), "交付清单", [
        "1. 软件功能需求分析文档",
        "2. 软件功能设计文档",
        "3. 软件产品说明书",
        "4. 软件功能测试报告",
        "5. 软件安装包及部署文档（含安装包压缩文件）",
        "6. 软件源文件压缩包",
        "7. 软件功能演示 PPT 文档",
    ])

    prs.save(path)


def write_delivery_index(path: Path, install_stats: PackageStats, source_stats: PackageStats) -> None:
    content = f"""# 交付说明

本目录为 2026-06-29 课程验收整理交付物。整理原则如下：

- 优先复用仓库现有 README、设计期文档、测试记录、部署文档和构建产物。
- 仓库中缺失的正式成品文档（如产品说明书、统一测试报告、正式 PPT）已补充生成。
- 当前仓库就是 GitHub 项目的本地工作副本，因此“GitHub 上已有内容”已按现有主线材料直接复用。

## 交付清单

1. `01-软件功能需求分析文档.docx`
2. `02-软件功能设计文档.docx`
3. `03-软件产品说明书.docx`
4. `04-软件功能测试报告.docx`
5. `05-软件安装包.zip`
6. `05-软件安装包及部署文档.docx`
7. `06-软件源文件.zip`
8. `07-软件功能演示PPT文档.pptx`

## 压缩包信息

- 安装包：`{install_stats.zip_name}`，约 {install_stats.zip_size_mb} MB，包含 {install_stats.file_count} 个文件。
- 源码包：`{source_stats.zip_name}`，约 {source_stats.zip_size_mb} MB，包含 {source_stats.file_count} 个文件。

## 说明

- 默认演示账号：`{DEFAULT_ACCOUNT}`
- 推荐演示路径：登录 -> 作业审批 -> 附件预览/下载 -> AI 评价 -> 教师复核 -> 统计/Excel 导出
- LoongArch 麒麟环境部署建议优先查看 `05-软件安装包及部署文档.docx`
"""
    path.write_text(content, encoding="utf-8")


def prepare_install_package() -> PackageStats:
    ensure_clean_dir(INSTALL_STAGE_DIR)
    INSTALL_APP_DIR.mkdir(parents=True, exist_ok=True)

    backend_dir = INSTALL_APP_DIR / "backend"
    backend_dir.mkdir(parents=True, exist_ok=True)
    shutil.copy2(BACKEND_JAR, backend_dir / BACKEND_JAR.name)

    frontend_dir = INSTALL_APP_DIR / "frontend-dist"
    copy_tree(FRONTEND_DIST, frontend_dir)

    ai_dir = INSTALL_APP_DIR / "ai-service"
    copy_tree(
        ROOT / "ai-service",
        ai_dir,
        ignore=shutil.ignore_patterns(".venv", "__pycache__", "*.pyc", ".DS_Store"),
    )
    remove_paths(ai_dir, [
        ".env",
        "logs",
        ".pytest_cache",
        ".mypy_cache",
        ".ruff_cache",
    ])

    deploy_dir = INSTALL_APP_DIR / "deployment"
    deploy_dir.mkdir(parents=True, exist_ok=True)
    for file_path in [
        ROOT / "README.md",
        ROOT / "pom.xml",
        ROOT / "docker-compose.yml",
        ROOT / "Dockerfile",
        ROOT / "start.sh",
        ROOT / "start.bat",
        ROOT / "start-kylin-loongarch-demo.sh",
        ROOT / "sql" / "schema.sql",
        ROOT / "docs" / "演示文档" / "麒麟LoongArch虚拟机部署指南.md",
    ]:
        shutil.copy2(file_path, deploy_dir / file_path.name)

    shutil.copy2(ROOT / "frontend" / "Dockerfile", deploy_dir / "frontend.Dockerfile")
    shutil.copy2(ROOT / "frontend" / "nginx.conf", deploy_dir / "frontend.nginx.conf")
    shutil.copy2(ROOT / "ai-service" / "Dockerfile", deploy_dir / "ai-service.Dockerfile")

    zip_path = DELIVERABLE_DIR / "05-软件安装包.zip"
    return zip_dir(INSTALL_STAGE_DIR, zip_path)


def prepare_source_package() -> PackageStats:
    ensure_clean_dir(SOURCE_STAGE_DIR)
    source_root = SOURCE_STAGE_DIR / "Teaching-System-source"
    copy_tree(
        ROOT,
        source_root,
        ignore=shutil.ignore_patterns(
            ".git",
            ".DS_Store",
            "node_modules",
            ".venv",
            "target",
            "output",
            "deliverables",
            "logs",
            "uploads",
            "__pycache__",
            "*.pyc",
        ),
    )
    remove_paths(source_root, [
        "ai-service/.env",
        "ai-service/logs",
        "frontend/dist",
        "out",
        "demo_script.docx",
        "tools",
    ])
    zip_path = DELIVERABLE_DIR / "06-软件源文件.zip"
    return zip_dir(SOURCE_STAGE_DIR, zip_path)


def generate_all() -> None:
    if not BACKEND_JAR.exists():
        raise FileNotFoundError(f"未找到后端打包产物：{BACKEND_JAR}")
    if not FRONTEND_DIST.exists():
        raise FileNotFoundError(f"未找到前端构建产物：{FRONTEND_DIST}")

    ensure_clean_dir(DELIVERABLE_DIR)
    ensure_clean_dir(BUILD_DIR)

    install_stats = prepare_install_package()
    source_stats = prepare_source_package()

    create_requirements_doc(DELIVERABLE_DIR / "01-软件功能需求分析文档.docx")
    create_design_doc(DELIVERABLE_DIR / "02-软件功能设计文档.docx")
    create_manual_doc(DELIVERABLE_DIR / "03-软件产品说明书.docx")
    create_test_doc(DELIVERABLE_DIR / "04-软件功能测试报告.docx")
    create_deployment_doc(DELIVERABLE_DIR / "05-软件安装包及部署文档.docx", install_stats, source_stats)
    create_ppt(DELIVERABLE_DIR / "07-软件功能演示PPT文档.pptx")
    write_delivery_index(DELIVERABLE_DIR / "00-交付说明.md", install_stats, source_stats)

    summary = {
        "deliverable_dir": str(DELIVERABLE_DIR),
        "files": sorted(p.name for p in DELIVERABLE_DIR.iterdir()),
    }
    (DELIVERABLE_DIR / "generation-summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    generate_all()
