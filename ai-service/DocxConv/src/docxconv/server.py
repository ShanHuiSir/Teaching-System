import json
import logging
import shutil
import tempfile
import time
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger(__name__)

from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

import config
from docxconv.converters.img import convert as convert_docx_to_images
from docxconv.converters.json import convert_obj
from evaluator.deepseek import evaluate as evaluate_content
from evaluator.deepseek import evaluate_stream

app = FastAPI(
    title="AI Evaluation Service",
    description="文件预处理与 AI 自动评分服务",
    version="0.8.0",
    servers=[
        {"url": config.SERVER_URL, "description": "当前环境"},
    ],
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=config.CORS_ORIGINS,
    allow_methods=["*"],
    allow_headers=["*"],
)

_rate_buckets: dict[str, tuple[int, int]] = {}


def _client_ip(request: Request) -> str:
    forwarded = request.headers.get("x-forwarded-for")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.client.host if request.client else "unknown"


def _check_ai_rate_limit(request: Request) -> None:
    limit = getattr(config, "AI_RATE_LIMIT_PER_MINUTE", 10)
    if limit <= 0:
        return
    key = _client_ip(request)
    window = int(time.time() // 60)
    bucket_window, count = _rate_buckets.get(key, (window, 0))
    if bucket_window != window:
        bucket_window, count = window, 0
    count += 1
    _rate_buckets[key] = (bucket_window, count)
    if count > limit:
        raise HTTPException(status_code=429, detail="AI接口调用过于频繁，请稍后再试")

# ── response models ─────────────────────────────────────────────────────

class PreprocessResponse(BaseModel):
    fileType: str = Field(description="文件类型：text / docx / pdf / xlsx / pptx / 图片后缀 / zip 等 / empty / unknown")
    originalFilename: str = Field(description="上传时的原始文件名")
    extractedText: str = Field(description="提取的纯文本内容，供 AI 评分使用")
    warnings: list[str] = Field(default_factory=list, description="预处理过程中的警告信息")
    renderStatus: str = Field("none", description="渲染状态：ok / degraded / failed / none")
    renderEngine: str = Field("none", description="渲染引擎：libreoffice / word / onlyoffice / none")
    renderWarnings: list[str] = Field(default_factory=list, description="渲染过程中的警告信息")
    structuredContent: Optional[list[Any]] = Field(None, description="仅 .docx 返回：标题/段落/列表/表格的结构化 JSON")

class EvaluateResponse(BaseModel):
    aiScore: float = Field(description="AI 评分，0-100")
    aiIssues: str = Field(description="扣分项列表，每条以 N. 开头，\\n 分隔")
    aiComment: str = Field(description="50-150 字综合评价")
    status: int = Field(description="固定为 1，表示正常返回")

class EvaluateRealResponse(PreprocessResponse):
    studentName: str = Field(description="学生姓名")
    aiScore: float = Field(description="AI 评分，0-100")
    aiIssues: str = Field(description="扣分项列表，每条以 N. 开头，\\n 分隔")
    aiComment: str = Field(description="50-150 字综合评价")
    dimensionScores: list[dict] = Field(default_factory=list, description="分维度评分明细，每项含 name/score/comment")
    status: int = Field(description="固定为 1，表示正常返回")

class FeatureStatus(BaseModel):
    available: bool = Field(description="该功能是否可用")
    detail: str = Field("", description="补充说明，不可用时描述原因")

class HealthResponse(BaseModel):
    status: str = Field(description="ok = 全部就绪 / degraded = 部分可选功能不可用 / unavailable = 核心功能不可用")
    version: str = Field(description="服务版本号")
    modules: dict[str, str] = Field(description="各子模块及其版本")
    features: dict[str, FeatureStatus] = Field(description="各功能/依赖的可用状态")

class DeepSeekHealthResponse(BaseModel):
    connected: bool = Field(description="是否成功连接到 DeepSeek API")
    model: str = Field(description="当前使用的模型")
    latencyMs: int = Field(0, description="API 调用耗时（毫秒），未连接时为 0")
    detail: str = Field("", description="补充说明")

# ── supported file types ──────────────────────────────────────────────

def _file_type(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in config.TEXT_EXTENSIONS:
        return "text"
    if suffix in config.TOOL_EXTENSIONS:
        return suffix.lstrip(".")
    return "unknown"

# ── extractors ────────────────────────────────────────────────────────

def _extract_text_file(content: bytes) -> str:
    for encoding in ("utf-8", "gbk", "utf-16"):
        try:
            return content.decode(encoding)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return content.decode("utf-8", errors="replace")

def _with_temp(content: bytes, suffix: str, fn):
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = Path(tmp.name)
    try:
        return fn(tmp_path)
    finally:
        tmp_path.unlink(missing_ok=True)

def _extract_docx_text(content: bytes) -> tuple[str, list]:
    def _do(tmp_path):
        structured = convert_obj(tmp_path)
        lines = []
        for item in structured:
            _flatten_content(item, lines)
        return "\n\n".join(lines), structured
    return _with_temp(content, ".docx", _do)

def _extract_docx_images_ocr(content: bytes) -> list[dict]:
    """Extract embedded images from a DOCX file and run OCR on each.

    Returns a list of {"context": str|None, "ocr_text": str} dicts.
    context is the surrounding paragraph text used for position matching.
    """
    import io as _io
    from PIL import Image as PILImage
    from docxconv.converters.extract_images import extract_images
    from screenshotproc.ocr import ocr_image

    def _do(tmp_path):
        images = extract_images(tmp_path)
        results = []
        for img_info in images:
            ocr_text = ""
            try:
                pil_img = PILImage.open(_io.BytesIO(img_info.raw_bytes))
                ocr_results = ocr_image(pil_img)
                if ocr_results:
                    ocr_results.sort(key=lambda r: (round(r["bbox"][1] / 30) * 30, r["bbox"][0]))
                    lines = [r["text"] for r in ocr_results if r["text"].strip()]
                    if lines:
                        ocr_text = "\n".join(lines)
            except Exception:
                pass
            if ocr_text:
                results.append({
                    "context": img_info.context,
                    "ocr_text": f"[图片 {img_info.index} ({img_info.width}x{img_info.height})]\n{ocr_text}",
                })
        return results

    return _with_temp(content, ".docx", _do)


def _insert_ocr_at_context(text: str, context: Optional[str], ocr_block: str) -> str:
    """Insert ocr_block into text near the matching context string.

    If context is given and found in text, inserts right after the last
    occurrence. Otherwise appends at the end.
    """
    if context:
        idx = text.rfind(context.strip())
        if idx != -1:
            insert_at = idx + len(context.strip())
            return text[:insert_at] + "\n\n" + ocr_block + "\n" + text[insert_at:]
    return text + "\n\n" + ocr_block


def _probe_docx_render(content: bytes) -> tuple[str, str, list[str]]:
    def _do(tmp_path):
        output_dir = tmp_path.parent / f"{tmp_path.stem}_render"
        try:
            convert_docx_to_images(tmp_path, output_dir)
            return "degraded", "libreoffice", ["DOCX 由 LibreOffice 渲染，复杂排版可能与 Word 存在差异"]
        finally:
            shutil.rmtree(output_dir, ignore_errors=True)

    try:
        return _with_temp(content, ".docx", _do)
    except Exception as e:
        return "failed", "libreoffice", [f"DOCX 渲染失败，已继续使用结构化文本进行预处理: {e}"]


def _flatten_content(item: dict, lines: list):
    if "heading" in item:
        prefix = "#" * item.get("level", 1)
        lines.append(f"{prefix} {item['heading']}")
        for sub in item.get("content", []):
            _flatten_content(sub, lines)
        for child in item.get("children", []):
            _flatten_content(child, lines)
    elif item.get("type") == "paragraph":
        lines.append(item.get("text", ""))
    elif item.get("type") == "list":
        for li in item.get("items", []):
            prefix = "- " if item.get("ordered") else "* "
            lines.append(f"{prefix}{li.get('text', '')}")
            for child in li.get("children", []):
                lines.append(f"  - {child.get('text', '')}")
    elif item.get("type") == "table":
        for row in item.get("rows", []):
            lines.append(" | ".join(row))

def _extract_pdf_text(content: bytes) -> str:
    import fitz
    def _do(tmp_path):
        doc = fitz.open(str(tmp_path))
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n\n".join(p for p in pages if p.strip())
    return _with_temp(content, ".pdf", _do)

def _extract_xlsx_text(content: bytes) -> str:
    import openpyxl
    def _do(tmp_path):
        wb = openpyxl.load_workbook(tmp_path, data_only=True, read_only=True)
        parts = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"--- 工作表: {name} ---\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    return _with_temp(content, ".xlsx", _do)

def _extract_pptx_text(content: bytes) -> str:
    from pptx import Presentation
    def _do(tmp_path):
        prs = Presentation(str(tmp_path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(c for c in cells if c))
            if texts:
                parts.append(f"--- 幻灯片 {i} ---\n" + "\n".join(texts))
        return "\n\n".join(parts)
    return _with_temp(content, ".pptx", _do)

def _extract_image_text(content: bytes) -> str:
    from PIL import Image
    import io
    from screenshotproc.ocr import ocr_image

    img = Image.open(io.BytesIO(content))
    results = ocr_image(img)
    if not results:
        return ""
    results.sort(key=lambda r: (round(r["bbox"][1] / 30) * 30, r["bbox"][0]))
    return "\n".join(r["text"] for r in results)

def _extract_archive_text(content: bytes, filename: str) -> tuple[str, list[str]]:
    from archiveproc.extractors import process

    tmpdir = Path(tempfile.mkdtemp())
    try:
        arc_path = tmpdir / filename
        arc_path.write_bytes(content)
        result = process(arc_path)
        texts: list[str] = []
        warnings: list[str] = []

        for entry in result.get("files", []):
            name = entry.get("path", "")
            ftype = entry.get("type", "")
            suffix = Path(name).suffix.lower()

            if ftype == "text":
                text = entry.get("text", "")
                if text:
                    texts.append(f"--- {name} ---\n{text}")
                elif entry.get("text_truncated"):
                    warnings.append(f"{name}: 文件过大或非 UTF-8 编码，已跳过")
            else:
                raw = entry.get("bytes")
                if raw is None:
                    continue
                if suffix == ".docx":
                    try:
                        t, _ = _extract_docx_text(raw)
                        if t:
                            texts.append(f"--- {name} ---\n{t}")
                    except Exception:
                        warnings.append(f"{name}: DOCX 解析失败")
                elif suffix in (".xlsx", ".xls"):
                    try:
                        t = _extract_xlsx_text(raw)
                        if t:
                            texts.append(f"--- {name} ---\n{t}")
                    except Exception:
                        warnings.append(f"{name}: 表格解析失败")
                elif suffix in (".pptx", ".ppt"):
                    try:
                        t = _extract_pptx_text(raw)
                        if t:
                            texts.append(f"--- {name} ---\n{t}")
                    except Exception:
                        warnings.append(f"{name}: 演示文稿解析失败")
                elif suffix == ".pdf":
                    try:
                        t = _extract_pdf_text(raw)
                        if t:
                            texts.append(f"--- {name} ---\n{t}")
                    except Exception:
                        warnings.append(f"{name}: PDF 解析失败")
                elif suffix in (".png", ".jpg", ".jpeg", ".bmp", ".webp"):
                    if not config.OCR_ENABLED:
                        continue
                    try:
                        t = _extract_image_text(raw)
                        if t:
                            texts.append(f"--- {name} ---\n{t}")
                    except Exception:
                        warnings.append(f"{name}: 图片 OCR 失败")

        return "\n\n".join(texts), warnings
    except Exception as e:
        return "", [f"解压失败: {e}"]
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)

def _extract_text(content: bytes, filename: str) -> tuple[str, list]:
    suffix = Path(filename).suffix.lower()
    if suffix in config.TEXT_EXTENSIONS:
        return _extract_text_file(content), []
    if suffix == ".docx":
        text, _ = _extract_docx_text(content)
        return text, []
    if suffix == ".doc":
        try:
            text, _ = _extract_docx_text(content)
            return text, ["旧版 .doc 格式，部分内容可能丢失"]
        except Exception:
            return "", [".doc 格式不受支持，请转为 .docx 后重试"]
    if suffix == ".pdf":
        text = _extract_pdf_text(content)
        if not text.strip():
            return "", ["PDF 中未提取到文本（可能是扫描件，当前版本不支持扫描件 OCR）"]
        return text, []
    if suffix in (".xlsx", ".xls"):
        try:
            return _extract_xlsx_text(content), []
        except Exception as e:
            return "", [f"表格解析失败: {e}"]
    if suffix in (".pptx", ".ppt"):
        try:
            return _extract_pptx_text(content), []
        except Exception as e:
            return "", [f"演示文稿解析失败: {e}"]
    if suffix in (".png", ".jpg", ".jpeg", ".bmp", ".webp"):
        if not config.OCR_ENABLED:
            return "", ["OCR 功能已禁用，无法识别图片中的文字"]
        try:
            text = _extract_image_text(content)
            if not text.strip():
                return "", ["图片中未识别到文字"]
            return text, ["图片内容由 OCR 识别，可能存在错字"]
        except Exception as e:
            return "", [f"OCR 识别失败: {e}"]
    if suffix in (".zip", ".tar", ".tar.xz", ".rar", ".7z"):
        text, warnings = _extract_archive_text(content, filename)
        return text, warnings
    try:
        return content.decode("utf-8"), ["无法识别文件类型，已尝试当作文本读取，结果可能不正确"]
    except UnicodeDecodeError:
        return "", [f"不支持的文件格式（{suffix}），无法提取文本"]

# ── /api/health ───────────────────────────────────────────────────────

@app.get(
    "/api/health",
    response_model=HealthResponse,
    summary="健康检查",
    description="返回服务运行状态、版本信息及各功能/依赖的可用状态。status 取值：ok（全部就绪）/ degraded（部分可选功能不可用）/ unavailable（核心功能不可用）。",
    tags=["系统"],
)
async def health():
    # ── core dependencies ────────────────────────────────────────────
    libreoffice_ok = bool(config.SOFFICE_PATH) and Path(config.SOFFICE_PATH).exists()
    pandoc_ok = bool(config.PANDOC_PATH)
    deepseek_ok = bool(config.DEEPSEEK_API_KEY)

    # ── optional dependencies ────────────────────────────────────────
    easyocr_ok = False
    try:
        import easyocr  # noqa: F401
        easyocr_ok = True
    except ImportError:
        pass

    archive_7z_ok = False
    try:
        import py7zr  # noqa: F401
        archive_7z_ok = True
    except ImportError:
        pass

    archive_rar_ok = False
    try:
        import rarfile  # noqa: F401
        archive_rar_ok = bool(rarfile.tool_setup())
    except Exception:
        pass

    # ── status ───────────────────────────────────────────────────────
    # DeepSeek/OCR/archive backends are feature readiness signals, not process
    # liveness. The service can still preprocess text and return fallback AI
    # results when those optional capabilities are unavailable.
    core_ok = True
    all_ok = core_ok and easyocr_ok and archive_7z_ok and archive_rar_ok
    if all_ok:
        status = "ok"
    elif core_ok:
        status = "degraded"
    else:
        status = "unavailable"

    return {
        "status": status,
        "version": "0.8.0",
        "modules": {
            "docxconv": "0.5.0",
            "screenshotproc": "0.1.0",
            "archiveproc": "0.1.0",
            "evaluator": "0.1.0",
        },
        "features": {
            "libreoffice": {
                "available": libreoffice_ok,
                "detail": "" if libreoffice_ok else "LibreOffice 未安装或路径不可用，DOCX 页面渲染不可用（结构化提取仍可用）",
            },
            "pandoc": {
                "available": pandoc_ok,
                "detail": "" if pandoc_ok else "Pandoc 未安装，备用渲染管线不可用",
            },
            "deepseek": {
                "available": deepseek_ok,
                "detail": "" if deepseek_ok else "DEEPSEEK_API_KEY 未配置，AI 评分不可用",
            },
            "easyocr": {
                "available": easyocr_ok,
                "detail": "" if easyocr_ok else "EasyOCR 未安装，截图 OCR 不可用",
            },
            "archive_7z": {
                "available": archive_7z_ok,
                "detail": "" if archive_7z_ok else "py7zr 未安装，7z 压缩包解压不可用",
            },
            "archive_rar": {
                "available": archive_rar_ok,
                "detail": "" if archive_rar_ok else "rarfile 或 unrar 未安装，RAR 解压不可用",
            },
        },
    }

# ── /api/health/deepseek ──────────────────────────────────────────────

@app.get(
    "/api/health/deepseek",
    response_model=DeepSeekHealthResponse,
    summary="DeepSeek 连通性深度检查",
    description="发起真实的 DeepSeek API 调用（max_tokens=1），测量连通性与延迟。与 /api/health 不同，此端点会消耗 API 配额。",
    tags=["系统"],
)
async def health_deepseek():
    if not config.DEEPSEEK_API_KEY:
        return {
            "connected": False,
            "model": config.DEEPSEEK_MODEL,
            "latencyMs": 0,
            "detail": "DEEPSEEK_API_KEY 未配置",
        }

    from evaluator.deepseek import _get_client

    client = _get_client()
    start = time.perf_counter()
    try:
        client.chat.completions.create(
            model=config.DEEPSEEK_MODEL,
            messages=[{"role": "user", "content": "ping"}],
            max_tokens=1,
            temperature=0,
        )
        elapsed = round((time.perf_counter() - start) * 1000)
        return {
            "connected": True,
            "model": config.DEEPSEEK_MODEL,
            "latencyMs": elapsed,
            "detail": "",
        }
    except Exception as e:
        elapsed = round((time.perf_counter() - start) * 1000)
        return {
            "connected": False,
            "model": config.DEEPSEEK_MODEL,
            "latencyMs": elapsed,
            "detail": f"API 调用失败: {e}",
        }

# ── /api/preprocess ───────────────────────────────────────────────────

MAX_UPLOAD_SIZE_MB = 200

@app.post(
    "/api/preprocess",
    response_model=PreprocessResponse,
    summary="文件预处理",
    description="上传任意格式文件，提取纯文本内容。支持 .txt/.docx/.pdf/.xlsx/.pptx/图片(OCR)/压缩包等。.docx 额外返回结构化 JSON。",
    tags=["预处理"],
)
async def preprocess(file: UploadFile = File(..., description="要预处理的文件")):
    if not file.filename:
        raise HTTPException(400, "No filename provided")
    if file.size and file.size > MAX_UPLOAD_SIZE_MB * 1024 * 1024:
        raise HTTPException(413, f"文件过大（最大 {MAX_UPLOAD_SIZE_MB} MB）")
    try:
        content = await file.read()
    except Exception:
        raise HTTPException(400, "Failed to read uploaded file")
    if not content:
        return {
            "fileType": "empty",
            "originalFilename": file.filename,
            "extractedText": "",
            "warnings": ["文件内容为空"],
            "renderStatus": "none",
            "renderEngine": "none",
            "renderWarnings": [],
        }
    ft = _file_type(file.filename)
    is_docx = Path(file.filename).suffix.lower() == ".docx"

    if is_docx:
        text, structured = _extract_docx_text(content)
        warnings = []
    else:
        text, warnings = _extract_text(content, file.filename)
        structured = None

    result = {
        "fileType": ft,
        "originalFilename": file.filename,
        "extractedText": text,
        "warnings": warnings,
        "renderStatus": "none",
        "renderEngine": "none",
        "renderWarnings": [],
    }
    if structured is not None:
        result["structuredContent"] = structured

    if is_docx:
        render_status, render_engine, render_warnings = _probe_docx_render(content)
        result["renderStatus"] = render_status
        result["renderEngine"] = render_engine
        result["renderWarnings"] = render_warnings
    if is_docx and config.OCR_ENABLED:
        try:
            ocr_results = _extract_docx_images_ocr(content)
            if ocr_results:
                text = result["extractedText"]
                for img in ocr_results:
                    text = _insert_ocr_at_context(text, img["context"], img["ocr_text"])
                result["extractedText"] = text
                warnings.append("文档中的图片内容由 OCR 识别，可能存在错字")
        except Exception:
            pass
    return result

# ── /api/evaluate (fake, backward compat) ──────────────────────────────

@app.post(
    "/api/evaluate",
    response_model=EvaluateResponse,
    summary="AI 评分（假）",
    description="不依赖 Python 预处理，直接返回硬编码评分。用于离线演示和开发调试。",
    tags=["评分"],
)
async def evaluate(
    studentName: str = "",
    fileName: str = "",
):
    if not studentName.strip():
        raise HTTPException(400, "studentName is required")
    return {
        "aiScore": 82.50,
        "aiIssues": (
            "1. 结构不够清晰，建议优化段落层次\n"
            "2. 缺少核心论点支撑材料\n"
            "3. 格式规范性不足，标题层级需统一"
        ),
        "aiComment": "整体完成度较好，但在结构组织上还有提升空间，建议加强逻辑连贯性。",
        "status": 1,
    }

# ── /api/evaluate/real (file upload → preprocess → AI) ────────────────

@app.post(
    "/api/evaluate/real",
    response_model=EvaluateRealResponse,
    summary="AI 评分（真实）",
    description=(
        "上传文件 → 预处理提取文本 → DeepSeek 评分。"
        "可选参数 subjectType（code/document/design/general）启用学科感知评分，"
        "可选参数 rubric（JSON 字符串）自定义评分维度。"
    ),
    tags=["评分"],
)
async def evaluate_real(
    request: Request,
    file: UploadFile = File(..., description="学生提交的作业文件"),
    studentName: str = Form(""),
    rubric: str = Form(None),
    subjectType: str = Form("general"),
):
    _check_ai_rate_limit(request)
    if not studentName.strip():
        raise HTTPException(400, "studentName is required")

    # Parse rubric JSON if provided
    rubric_dict = None
    if rubric and rubric.strip():
        try:
            rubric_dict = json.loads(rubric)
        except json.JSONDecodeError:
            raise HTTPException(422, "Invalid rubric JSON: 无法解析为 JSON")

    # Validate rubric early
    if rubric_dict is not None:
        from evaluator.deepseek import _validate_rubric
        try:
            _validate_rubric(rubric_dict.get("dimensions", []))
        except ValueError as e:
            raise HTTPException(422, f"Invalid rubric: {e}")

    preprocess_result = await preprocess(file)
    try:
        eval_result = evaluate_content(
            preprocess_result["extractedText"],
            studentName,
            rubric_json=rubric_dict,
            subject_type=subjectType,
        )
    except RuntimeError as e:
        logger.warning("AI evaluation failed, falling back to default scores: %s", e)
        preprocess_result["warnings"].append(f"DeepSeek 评分失败，已降级为默认评分: {e}")
        preprocess_result["aiScore"] = 82.50
        preprocess_result["aiIssues"] = (
            "1. AI 评分服务暂时不可用，以下为默认提示\n"
            "2. 请稍后重试或联系教师人工评阅"
        )
        preprocess_result["aiComment"] = "AI 评分服务暂时不可用，当前分数为系统默认值（不代表真实评价）。请稍后重试或联系教师。"
        preprocess_result["dimensionScores"] = []
        preprocess_result["status"] = 1
        preprocess_result["studentName"] = studentName
        return preprocess_result
    return {
        "studentName": studentName,
        "originalFilename": preprocess_result["originalFilename"],
        "fileType": preprocess_result["fileType"],
        "extractedText": preprocess_result["extractedText"],
        "warnings": preprocess_result["warnings"],
        "renderStatus": preprocess_result["renderStatus"],
        "renderEngine": preprocess_result["renderEngine"],
        "renderWarnings": preprocess_result["renderWarnings"],
        "aiScore": eval_result["aiScore"],
        "aiIssues": eval_result["aiIssues"],
        "aiComment": eval_result["aiComment"],
        "dimensionScores": eval_result.get("dimensionScores", []),
        "status": 1,
    }

# ── /api/evaluate/stream (SSE) ────────────────────────────────────────────

@app.post(
    "/api/evaluate/stream",
    summary="AI 评分（流式 SSE）",
    description=(
        "上传文件 → 预处理提取文本 → DeepSeek 流式评分。"
        "以 SSE (Server-Sent Events) 推送思考过程和评分结果。"
        "事件类型：start / reasoning / content / result / error / done。"
    ),
    tags=["评分"],
)
async def evaluate_stream_endpoint(
    request: Request,
    file: UploadFile = File(..., description="学生提交的作业文件"),
    studentName: str = Form(""),
    rubric: str = Form(None),
    subjectType: str = Form("general"),
):
    _check_ai_rate_limit(request)
    if not studentName.strip():
        raise HTTPException(400, "studentName is required")

    rubric_dict = None
    if rubric and rubric.strip():
        try:
            rubric_dict = json.loads(rubric)
        except json.JSONDecodeError:
            raise HTTPException(422, "Invalid rubric JSON: 无法解析为 JSON")
        try:
            from evaluator.deepseek import _validate_rubric
            _validate_rubric(rubric_dict.get("dimensions", []))
        except ValueError as e:
            raise HTTPException(422, f"Invalid rubric: {e}")

    preprocess_result = await preprocess(file)

    def _sse_generator():
        """Sync generator wrapping evaluate_stream for StreamingResponse."""
        try:
            for sse_frame in evaluate_stream(
                preprocess_result["extractedText"],
                studentName,
                rubric_json=rubric_dict,
                subject_type=subjectType,
            ):
                yield sse_frame
        except Exception:
            logger.exception("Unhandled error in SSE stream")
            yield f"event: error\ndata: {json.dumps({'message': '服务器内部错误', 'code': 'INTERNAL_ERROR'})}\n\n"
            yield f"event: done\ndata: {json.dumps({})}\n\n"

    return StreamingResponse(
        _sse_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


# ── Pydantic models for eval-log ────────────────────────────────────────

class EvalLogRequest(BaseModel):
    student_name: str = Field(..., description="学生姓名")
    ai_score: float = Field(..., description="AI 评分")
    ai_issues: str = Field("", description="AI 发现的问题")
    ai_comment: str = Field("", description="AI 综合评语")
    dimension_scores: list[dict] = Field(default_factory=list, description="AI 分维度评分")
    teacher_score: float = Field(..., description="教师最终评分")
    teacher_comment: str = Field(..., description="教师最终评语")


# ── /api/eval-log ───────────────────────────────────────────────────────

@app.post(
    "/api/eval-log",
    response_model=dict,
    summary="写入最终评价日志",
    description="由 Java 后端在教师批改完成后调用，将 AI + 教师评价合并写入 JSONL 日志文件。",
    tags=["评价日志"],
)
async def eval_log(body: EvalLogRequest):
    from evaluator.logger import log_final_evaluation

    log_final_evaluation(
        student_name=body.student_name,
        ai_score=body.ai_score,
        ai_issues=body.ai_issues,
        ai_comment=body.ai_comment,
        dimension_scores=body.dimension_scores,
        teacher_score=body.teacher_score,
        teacher_comment=body.teacher_comment,
    )
    return {"status": "ok"}


def main():
    import uvicorn
    uvicorn.run("docxconv.server:app", host="0.0.0.0", port=8000, reload=False)

if __name__ == "__main__":
    main()
